import google.generativeai as genai
import json
from pathlib import Path
import subprocess
import platform
import logging
import cv2
from typing import Dict

logger = logging.getLogger(__name__)


class FileExtractor:
    """Enhanced file extractor with combined content extraction and prototype detection"""
    
    def __init__(self, api_key: str, model: str = "gemini-2.0-flash-exp"):
        genai.configure(api_key=api_key)
        self.model = model
        self._client = None
    
    def _get_client(self):
        """Lazy load Gemini client"""
        if self._client is None:
            self._client = genai.GenerativeModel(model_name=self.model)
        return self._client
    
    def extract_content(self, file_path: Path) -> Dict[str, str]:
        """
        Extract text content AND determine content type in single API call
        
        Returns: {
            "content": "extracted text...",
            "content_type": "Prototype" or "Text"
        }
        """
        
        logger.info(f"Extracting file: {file_path}")
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        extension = file_path.suffix.lower()
        
        # Handle video files with frame extraction
        if extension in ['.mp4', '.mov', '.avi', '.mkv']:
            return self._extract_video_with_frames(file_path)
        
        # Handle other files
        file_path, mime_type = self._prepare_file(file_path)
        
        with open(file_path, 'rb') as f:
            file_data = f.read()
        
        # ✅ Combined prompt: extraction + classification in one call
        prompt = self._create_combined_prompt(file_path.name)
        
        try:
            client = self._get_client()
            # Upload file to Gemini
            uploaded_file = genai.upload_file(file_path, mime_type=mime_type)
            response = client.generate_content([uploaded_file, prompt])
            
            response_text = response.text.strip()
            
            # Clean markdown
            if response_text.startswith("```json"):
                response_text = response_text.replace("```json", "").replace("```", "")
            elif response_text.startswith("```"):
                response_text = response_text.replace("```", "")
            
            # Parse JSON response
            result = json.loads(response_text)
            
            extracted_text = result.get('content', '')
            content_type = result.get('content_type', 'Text')
            
            logger.info(f"✓ Extracted {len(extracted_text)} characters | Type: {content_type}")
            
            return {
                'content': extracted_text,
                'content_type': content_type
            }
            
        except Exception as e:
            logger.error(f"✗ Extraction failed: {e}")
            return {
                'content': f"Extraction failed: {str(e)}",
                'content_type': 'Text'
            }
    
    def _extract_video_with_frames(self, video_path: Path) -> Dict[str, str]:
        """Extract frames from video and analyze with single API call"""
        
        logger.info(f"  📹 Extracting frames from video every 10 seconds...")
        
        frames = self._extract_frames_every_n_seconds(video_path, interval_seconds=10)
        
        if not frames:
            return {
                'content': "Video analysis failed: No frames extracted",
                'content_type': 'Text'
            }
        
        logger.info(f"  ✓ Extracted {len(frames)} frames from video")
        
        # Save frames temporarily and upload
        import tempfile
        frame_files = []
        for idx, frame_data in enumerate(frames):
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
            temp_file.write(frame_data)
            temp_file.close()
            uploaded = genai.upload_file(temp_file.name, mime_type="image/jpeg")
            frame_files.append(uploaded)
        
        # ✅ Combined prompt for video
        prompt = f"""
Analyze this video presentation by examining {len(frames)} frames extracted every 10 seconds.

Video: {video_path.name}

Provide:
1. Comprehensive content extraction (main topic, key points, visual elements, technical details)
2. Classification as "Prototype" or "Text"

CLASSIFICATION RULES:
- **Prototype**: Prototype screenshots, prototype working demo evidence.
- **Text**: Only concepts, no technical implementation, theoretical only

Return JSON:
{{
  "content": "Detailed text summary of all information in the video",
  "content_type": "Prototype" or "Text"
}}
"""
        
        contents = frame_files + [prompt]
        
        try:
            client = self._get_client()
            response = client.generate_content(contents)
            
            # Clean up temp files
            import os
            for uploaded in frame_files:
                try:
                    temp_path = uploaded.name
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
                except:
                    pass
            
            response_text = response.text.strip()
            if response_text.startswith("```json"):
                response_text = response_text.replace("```json", "").replace("```", "")
            
            result = json.loads(response_text)
            
            logger.info(f"  ✓ Video analysis: {len(result['content'])} chars | Type: {result['content_type']}")
            
            return result
            
        except Exception as e:
            logger.error(f"  ✗ Video analysis failed: {e}")
            return {
                'content': f"Video analysis failed: {str(e)}",
                'content_type': 'Text'
            }
    
    def _extract_frames_every_n_seconds(self, video_path: Path, interval_seconds: int = 10) -> list:
        """Extract frames from video every N seconds using OpenCV"""
        frames = []
        
        try:
            cap = cv2.VideoCapture(str(video_path))
            
            if not cap.isOpened():
                logger.error(f"Could not open video: {video_path}")
                return frames
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = total_frames / fps if fps > 0 else 0
            
            logger.info(f"    Video: {duration:.1f}s, {fps:.1f} FPS")
            
            frame_interval = int(fps * interval_seconds)
            frame_count = 0
            extracted_count = 0
            
            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                
                if frame_count % frame_interval == 0:
                    success, buffer = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
                    if success:
                        frames.append(buffer.tobytes())
                        extracted_count += 1
                        timestamp = frame_count / fps
                        logger.info(f"    Frame {extracted_count} @ {timestamp:.1f}s")
                
                frame_count += 1
            
            cap.release()
            
        except Exception as e:
            logger.error(f"  ✗ Frame extraction failed: {e}")
        
        return frames
    
    def _create_combined_prompt(self, filename: str) -> str:
        return f"""
CRITICAL: You MUST extract detailed content AND classify the type.

File: {filename}

TASK 1 - CONTENT EXTRACTION (MANDATORY - DO NOT SKIP):
Extract ALL information in great detail (minimum 500 characters):
- All text content word-by-word
- Technical details and specifications
- Visual elements (describe diagrams, charts, images)
- Architecture and design information
- Code snippets, data, metrics

TASK 2 - CLASSIFICATION:
Classify as "Prototype" or "Text" based on technical implementation evidence.

**PROTOTYPE indicators**: Prototype screenshots, Prototype working demo
**TEXT indicators**: Only concepts, no implementation details

MANDATORY OUTPUT - Return valid JSON with BOTH fields:
{{
  "content": "DETAILED comprehensive extraction of ALL content from this file (minimum 500 characters)",
  "content_type": "Prototype" or "Text"
}}

CRITICAL: The "content" field MUST contain the full detailed extraction. Do not return empty or minimal content.
"""

    
    def _prepare_file(self, file_path: Path) -> tuple:
        """Convert unsupported formats to PDF"""
        extension = file_path.suffix.lower()
        mime_type = self._get_mime_type(extension)
        
        if extension in ['.pptx', '.docx']:
            logger.info(f"    🔄 Converting {extension.upper()} to PDF...")
            pdf_path = self._convert_to_pdf(file_path)
            return pdf_path, 'application/pdf'
        
        return file_path, mime_type
    
    def _convert_to_pdf(self, file_path: Path) -> Path:
        """Convert PPTX/DOCX to PDF using LibreOffice"""
        output_pdf = file_path.with_suffix('.converted.pdf')
        
        if output_pdf.exists():
            return output_pdf
        
        try:
            soffice_cmd = 'soffice' if platform.system() == 'Windows' else 'libreoffice'
            subprocess.run([
                soffice_cmd, '--headless', '--convert-to', 'pdf',
                '--outdir', str(file_path.parent), str(file_path)
            ], capture_output=True, timeout=60, check=True)
            
            converted_file = file_path.with_suffix('.pdf')
            if converted_file.exists():
                if output_pdf.exists():
                    output_pdf.unlink()
                converted_file.rename(output_pdf)
                return output_pdf
                
        except Exception as e:
            logger.warning(f"    ⚠ LibreOffice conversion failed: {e}")
        
        return self._fallback_conversion(file_path, output_pdf)
    
    def _fallback_conversion(self, file_path: Path, output_pdf: Path) -> Path:
        """Fallback conversion for PPTX/DOCX"""
        extension = file_path.suffix.lower()
        
        if extension == '.pptx':
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            prs = Presentation(file_path)
            c = canvas.Canvas(str(output_pdf), pagesize=letter)
            width, height = letter
            
            for slide_num, slide in enumerate(prs.slides, 1):
                c.setFont("Helvetica-Bold", 14)
                c.drawString(50, height - 50, f"Slide {slide_num}")
                y = height - 80
                c.setFont("Helvetica", 10)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        for line in shape.text.split('\n'):
                            if y < 100:
                                c.showPage()
                                y = height - 80
                            c.drawString(50, y, line[:100])
                            y -= 15
                c.showPage()
            c.save()
            return output_pdf
        
        elif extension == '.docx':
            from docx import Document
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            doc = Document(file_path)
            pdf = SimpleDocTemplate(str(output_pdf))
            styles = getSampleStyleSheet()
            story = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    story.append(Paragraph(para.text, styles['BodyText']))
                    story.append(Spacer(1, 12))
            
            pdf.build(story)
            return output_pdf
    
    def _get_mime_type(self, extension: str) -> str:
        """Get MIME type from extension"""
        mime_types = {
            '.pdf': 'application/pdf',
            '.mp4': 'video/mp4', '.mov': 'video/quicktime', '.avi': 'video/x-msvideo',
            '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', '.png': 'image/png',
            '.webp': 'image/webp', '.gif': 'image/gif'
        }
        return mime_types.get(extension, 'application/octet-stream')
